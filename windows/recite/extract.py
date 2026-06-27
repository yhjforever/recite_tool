"""源文件文本抽取：支持 .pdf / .pptx / .txt / .md，带轻量去噪与缓存。"""
import re
import json
import logging
import hashlib
import warnings
from pathlib import Path

warnings.filterwarnings("ignore")
# pypdf 会用 logger 直接打印 "Ignoring wrong pointing object ..."，压到 ERROR 级别
logging.getLogger("pypdf").setLevel(logging.ERROR)

SUPPORTED = {".pdf", ".pptx", ".txt", ".md"}

# 课件常见噪声（整行命中即删）。系统提示词里也会让模型再清一遍。
NOISE_SUBSTR = [
    "图片来自互联网", "仅供学习参考", "仅供参考",
]
# 整行完全等于这些时删除（页眉/机构名）。保持学科无关，默认空；
# 真正的页眉清理交由系统提示词的"去噪声"铁律完成。
NOISE_EQUAL: set[str] = set()

HEADING_RE = re.compile(
    r"^(第[一二三四五六七八九十百0-9]+[章节篇]"
    r"|[一二三四五六七八九十]+、"
    r"|（[一二三四五六七八九十0-9]+）"
    r"|\([一二三四五六七八九十0-9]+\))"
)


def _clean_lines(raw: str, keep_digit_lines: bool = False) -> list[str]:
    out = []
    for ln in raw.split("\n"):
        s = ln.strip()
        if not s:
            continue
        # 纯数字行通常是页码，删之；但大纲是表格，PDF 抽取会把标题里的数字
        # （如 "Chapter30-\n35"、章次范围）拆成独立数字行，删了会破坏标题匹配，
        # 故抽取大纲时 keep_digit_lines=True 保留。
        if not keep_digit_lines and s.isdigit():
            continue
        if s in NOISE_EQUAL:
            continue
        if any(n in s for n in NOISE_SUBSTR):
            continue
        out.append(s)
    return out


# ---------- 各格式读取 ----------
def _read_pdf(path: Path) -> str:
    import pypdf
    r = pypdf.PdfReader(str(path))
    pages = []
    for p in r.pages:
        pages.append(p.extract_text() or "")
    return "\n".join(pages)


def _read_pptx(path: Path) -> str:
    try:
        from pptx import Presentation
    except ImportError:
        raise SystemExit("检测到 .pptx，但未安装 python-pptx：pip install python-pptx")
    prs = Presentation(str(path))
    chunks = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                chunks.append(shape.text_frame.text)
            if shape.has_table:
                for row in shape.table.rows:
                    chunks.append("\t".join(c.text for c in row.cells))
    return "\n".join(chunks)


def _read_plain(path: Path) -> str:
    """按 Windows 常见编码依次尝试：utf-8-sig → gb18030 → utf-16 → cp936。
    不用 errors='ignore' 静默吞字；都失败才宽松解码并校验乱码比例。"""
    raw = path.read_bytes()
    for enc in ("utf-8-sig", "gb18030", "utf-16"):
        try:
            return raw.decode(enc)
        except (UnicodeDecodeError, LookupError):
            continue
    text = raw.decode("utf-8", errors="replace")
    # 乱码（替换符 U+FFFD）比例过高 → 报错，避免拿乱码去喂 DeepSeek
    if text and text.count("�") / max(len(text), 1) > 0.02:
        raise SystemExit(f"文件疑似非 UTF-8/GBK 编码、解码出现大量乱码：{path.name}\n"
                         f"请用记事本『另存为 UTF-8』后重试。")
    return text


def read_raw(path: Path) -> str:
    ext = path.suffix.lower()
    if ext == ".pdf":
        return _read_pdf(path)
    if ext == ".pptx":
        return _read_pptx(path)
    if ext in (".txt", ".md"):
        return _read_plain(path)
    raise ValueError(f"不支持的文件类型: {path.name}")


# ---------- 带缓存的清洗文本 ----------
def _cache_key(path: Path) -> str:
    """用文件内容 hash 作 key：内容变了缓存必失效（不再依赖 size/秒级 mtime）。"""
    h = hashlib.md5()
    with open(path, "rb") as f:
        for chunk in iter(lambda: f.read(1 << 20), b""):
            h.update(chunk)
    return f"{path.stem}.{h.hexdigest()[:12]}.txt"


def extract_text(path: Path, cache_dir: Path, force: bool = False,
                 keep_digit_lines: bool = False) -> str:
    """抽取并清洗文本，结果按内容 hash 缓存到 cache_dir。
    keep_digit_lines=True 时保留纯数字行（用于大纲表格，避免标题数字被当页码删掉），
    并用独立缓存文件名，避免两种清洗模式互相污染。"""
    cache_dir.mkdir(parents=True, exist_ok=True)
    name = _cache_key(path)
    if keep_digit_lines:
        name = name[:-4] + ".kd.txt"
    cache_file = cache_dir / name
    if cache_file.exists() and not force:
        return cache_file.read_text(encoding="utf-8")
    text = "\n".join(_clean_lines(read_raw(path), keep_digit_lines))
    cache_file.write_text(text, encoding="utf-8")
    return text


def fingerprint(path: Path, cache_dir: Path) -> dict:
    """给审计步骤用的文件指纹：文件名 + 概览 + 疑似标题行。"""
    text = extract_text(path, cache_dir)
    lines = [l for l in text.split("\n") if l.strip()]
    headings = [l for l in lines if HEADING_RE.match(l)][:14]
    head = text[:400].replace("\n", " ")
    return {
        "file": path.name,
        "chars": len(text),
        "head": head,
        "headings": headings,
    }


# ---------- 文件分类 ----------
def is_supported(path: Path) -> bool:
    return path.suffix.lower() in SUPPORTED


def detect_outline(files: list[Path], explicit: str = "") -> Path | None:
    if explicit:
        for f in files:
            if f.name == explicit:
                return f
        raise SystemExit(f"outline_file 指定为 {explicit}，但母文件夹中找不到该文件")
    kw = ("大纲", "教学大纲", "syllabus", "outline")
    cands = [f for f in files if any(k in f.name.lower() or k in f.name for k in kw)]
    if len(cands) == 1:
        return cands[0]
    if len(cands) > 1:
        # 优先含“教学大纲”的
        best = [f for f in cands if "教学大纲" in f.name]
        return best[0] if best else cands[0]
    return None


def list_source_files(cfg) -> tuple[Path, list[Path]]:
    """返回 (大纲文件, 课件文件列表)。"""
    import fnmatch
    files = [f for f in sorted(cfg.source_dir.iterdir())
             if f.is_file() and is_supported(f)]
    outline = detect_outline(files, cfg.outline_file)
    if outline is None:
        raise SystemExit("未能自动识别大纲文件：请在 config.yaml 用 outline_file 指定其完整文件名。")

    def ignored(name: str) -> bool:
        low = name.lower()
        return any(fnmatch.fnmatch(low, pat.lower()) for pat in cfg.ignore)

    sources = [f for f in files if f != outline and not ignored(f.name)]
    return outline, sources
